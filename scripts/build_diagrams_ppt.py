from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

base = Path(r"c:\Users\Nico\Documents\NILA")
out_dir = base / "docs" / "diagrams" / "out"
out_ppt = base / "docs" / "diagrams" / "NILA-diagramas-logicos.pptx"

images = [
    out_dir / "01-contexto.png",
    out_dir / "02-backend-modular.png",
    out_dir / "03-auth-tenancy-secuencia.png",
    out_dir / "04-reserva-turno-secuencia.png",
    out_dir / "05-vista-objetivo.png",
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
sw, sh = prs.slide_width, prs.slide_height

for path in images:
    with Image.open(path) as im:
        wpx, hpx = im.size
    scale = min(sw / wpx, sh / hpx)
    pw, ph = int(wpx * scale), int(hpx * scale)
    left, top = int((sw - pw) / 2), int((sh - ph) / 2)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(path), left, top, width=pw, height=ph)

prs.save(out_ppt)
print(out_ppt)
