from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

base = Path(r"c:\Users\Nico\Documents\NILA")
web = base / "mockups" / "images" / "board-ready"
mobile = base / "mockups" / "images" / "board-ready-mobile"
out = base / "docs" / "NILA-investor-short-web-mobile.pptx"

slides = [
    web / "portada-ejecutiva.png",
    web / "cliente-busqueda.png",
    web / "cliente-dashboard.png",
    web / "cliente-lista-espera.png",
    web / "pro-dashboard.png",
    web / "pro-cancelaciones-automaticas.png",
    web / "pro-pos-facturacion.png",
    web / "pro-promocionar.png",
    web / "pro-analytics-avanzado.png",
    mobile / "mobile-cliente-home.png",
    mobile / "mobile-cliente-busqueda.png",
    mobile / "mobile-pro-home.png",
    mobile / "mobile-pro-pos.png",
    web / "accesibilidad-config.png",
]

for p in slides:
    if not p.exists():
        raise SystemExit(f"Missing image: {p}")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
sw = prs.slide_width
sh = prs.slide_height

for path in slides:
    with Image.open(path) as im:
        wpx, hpx = im.size
    scale = min(sw / wpx, sh / hpx)
    pw = int(wpx * scale)
    ph = int(hpx * scale)
    left = int((sw - pw) / 2)
    top = int((sh - ph) / 2)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(path), left, top, width=pw, height=ph)

prs.save(out)
print(f"Created: {out}")
print(f"Slides: {len(slides)}")
