from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

base = Path(r"c:\Users\Nico\Documents\NILA")
img_dir = base / "mockups" / "images" / "board-ready"
out = base / "docs" / "NILA-board-ready.pptx"

slides = [
    "portada-ejecutiva.png",
    "cliente-busqueda.png",
    "cliente-resultados.png",
    "cliente-login.png",
    "cliente-dashboard.png",
    "cliente-calendario.png",
    "cliente-rutina.png",
    "cliente-ia-postura.png",
    "cliente-pagos-facturas.png",
    "pro-dashboard.png",
    "pro-agenda.png",
    "pro-cancelaciones-automaticas.png",
    "pro-whatsapp-automatizaciones.png",
    "pro-ficha-cliente.png",
    "pro-pos-facturacion.png",
    "pro-membresias-suscripciones.png",
    "pro-promocionar.png",
    "pro-referidos.png",
    "pro-firma-digital.png",
    "pro-miniweb-publica.png",
    "pro-integraciones.png",
    "pro-analytics-avanzado.png",
    "accesibilidad-config.png",
]

missing = [name for name in slides if not (img_dir / name).exists()]
if missing:
    raise SystemExit(f"Missing images: {missing}")

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

sw = prs.slide_width
sh = prs.slide_height

for name in slides:
    path = img_dir / name
    with Image.open(path) as im:
        wpx, hpx = im.size

    # pixels are converted with same DPI factor, ratio remains valid
    scale = min(sw / wpx, sh / hpx)
    pic_w = int(wpx * scale)
    pic_h = int(hpx * scale)
    left = int((sw - pic_w) / 2)
    top = int((sh - pic_h) / 2)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(path), left, top, width=pic_w, height=pic_h)

prs.save(out)
print(f"Created: {out}")
print(f"Slides: {len(slides)}")
