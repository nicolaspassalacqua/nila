from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

base = Path(r"c:\Users\Nico\Documents\NILA")

configs = [
    {
        "web": base / "mockups" / "images" / "premium-dark-full",
        "mobile": base / "mockups" / "images" / "premium-dark-mobile",
        "out": base / "docs" / "NILA-investor-short-premium-dark-web-mobile.pptx",
    },
    {
        "web": base / "mockups" / "images" / "premium-wellness-full",
        "mobile": base / "mockups" / "images" / "premium-wellness-mobile",
        "out": base / "docs" / "NILA-investor-short-premium-wellness-web-mobile.pptx",
    },
]

slide_names = [
    ("web", "portada-ejecutiva.png"),
    ("web", "cliente-busqueda.png"),
    ("web", "cliente-dashboard.png"),
    ("web", "cliente-lista-espera.png"),
    ("web", "pro-dashboard.png"),
    ("web", "pro-cancelaciones-automaticas.png"),
    ("web", "pro-pos-facturacion.png"),
    ("web", "pro-promocionar.png"),
    ("web", "pro-analytics-avanzado.png"),
    ("mobile", "mobile-cliente-home.png"),
    ("mobile", "mobile-cliente-busqueda.png"),
    ("mobile", "mobile-pro-home.png"),
    ("mobile", "mobile-pro-pos.png"),
    ("web", "accesibilidad-config.png"),
]

for cfg in configs:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    sw = prs.slide_width
    sh = prs.slide_height

    for scope, name in slide_names:
        path = cfg[scope] / name
        if not path.exists():
            raise SystemExit(f"Missing image: {path}")

        with Image.open(path) as im:
            wpx, hpx = im.size

        scale = min(sw / wpx, sh / hpx)
        pw = int(wpx * scale)
        ph = int(hpx * scale)
        left = int((sw - pw) / 2)
        top = int((sh - ph) / 2)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(path), left, top, width=pw, height=ph)

    prs.save(cfg["out"])
    print(f"Created: {cfg['out']}")
    print("Slides: 14")
